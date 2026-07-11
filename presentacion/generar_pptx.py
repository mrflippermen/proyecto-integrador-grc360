# -*- coding: utf-8 -*-
"""Genera la presentación ejecutiva CyberRisk 360 en .pptx (tema oscuro)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- Paleta ----
BG     = RGBColor(0x0e,0x13,0x1f)
BG2    = RGBColor(0x12,0x1a,0x29)
SURF   = RGBColor(0x18,0x22,0x33)
SURF2  = RGBColor(0x1e,0x2a,0x3e)
BORDER = RGBColor(0x29,0x35,0x4b)
TEXT   = RGBColor(0xe8,0xed,0xf5)
WHITE  = RGBColor(0xff,0xff,0xff)
MUTED  = RGBColor(0x93,0xa1,0xb8)
DIM    = RGBColor(0x64,0x74,0x8b)
BRAND  = RGBColor(0x4c,0x82,0xf7)
BRAND2 = RGBColor(0x34,0xd3,0xc0)
BAJO   = RGBColor(0x2e,0x9e,0x5b)
MEDIO  = RGBColor(0xe0,0xa8,0x00)
ALTO   = RGBColor(0xe8,0x59,0x0c)
CRIT   = RGBColor(0xd6,0x3a,0x3a)

HEAD = "Segoe UI"
BODY = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height
ML = Inches(0.9)  # margen izquierdo


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def rect(s, l, t, w, h, fill=SURF, line=BORDER, radius=0.055, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = s.shapes.add_shape(shape, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    try:
        if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
            sp.adjustments[0] = radius
    except Exception:
        pass
    return sp


def _set(run, text, size, color, bold, font, italic=False):
    run.text = text; f = run.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = font


def txt(s, l, t, w, h, text, size, color=TEXT, bold=False, font=BODY,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, spacing=1.0):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = spacing
    _set(p.add_run(), text, size, color, bold, font, italic)
    return tb


def rich(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0):
    """runs: lista de (text,size,color,bold[,font])."""
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = spacing
    for r in runs:
        font = r[4] if len(r) > 4 else BODY
        _set(p.add_run(), r[0], r[1], r[2], r[3], font)
    return tb


def shape_text(sp, text, size, color, bold=True, font=HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=True):
    tf = sp.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = Pt(4); tf.margin_right = Pt(4); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]; p.alignment = align
    _set(p.add_run(), text, size, color, bold, font)


def eyebrow(s, text, t=Inches(0.85)):
    txt(s, ML, t, Inches(11), Inches(0.4), text.upper(), 12.5, BRAND2, bold=True, font=HEAD)


def title(s, text, t=Inches(1.25), color=TEXT, size=34):
    txt(s, ML, t, Inches(11.5), Inches(1.1), text, size, color, bold=True, font=HEAD, spacing=1.02)


def footer(s, section):
    # marca
    m = rect(s, ML, Inches(6.95), Inches(0.26), Inches(0.26), fill=BRAND, line=None, radius=0.3)
    txt(s, Inches(1.24), Inches(6.93), Inches(3), Inches(0.3), "CyberRisk 360", 10, MUTED, bold=True, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(8.5), Inches(6.93), Inches(3.9), Inches(0.3), section, 10, DIM, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def chips(s, items, l, t, on=True):
    x = l
    for it in items:
        w = Inches(0.34 + 0.108 * len(it))
        c = rect(s, x, t, w, Inches(0.44), fill=SURF2 if not on else RGBColor(0x1a,0x2a,0x46),
                 line=BORDER if not on else RGBColor(0x2d,0x4d,0x87), radius=0.28)
        shape_text(c, it, 11, MUTED if not on else RGBColor(0xa9,0xc4,0xfb), bold=True, font=HEAD, wrap=False)
        x = x + w + Inches(0.16)


def card(s, l, t, w, h, num=None, head=None, body=None, head_color=TEXT, fill=SURF):
    c = rect(s, l, t, w, h, fill=fill)
    pad = Inches(0.26)
    y = t + Inches(0.24)
    if num is not None:
        b = rect(s, l + pad, y, Inches(0.4), Inches(0.4), fill=BG2, line=BORDER, radius=0.22)
        shape_text(b, str(num), 13, BRAND, bold=True, font=HEAD)
        y = y + Inches(0.52)
    if head:
        txt(s, l + pad, y, w - pad*2, Inches(0.4), head, 15.5, head_color, bold=True, font=HEAD, spacing=1.0)
        y = y + Inches(0.46)
    if body:
        txt(s, l + pad, y, w - pad*2, t + h - y - Inches(0.14), body, 12, MUTED, spacing=1.08)
    return c


def bullets(s, l, t, w, items, size=15, gap=0.16):
    y = t
    for lead, rest in items:
        # punto
        d = rect(s, l, y + Inches(0.09), Inches(0.11), Inches(0.11), fill=BRAND, line=None, radius=0.5)
        runs = []
        if lead:
            runs.append((lead + " ", size, WHITE, True))
        if rest:
            runs.append((rest, size, MUTED, False))
        tb = rich(s, l + Inches(0.3), y, w - Inches(0.3), Inches(1.2), runs, spacing=1.12)
        # estimar alto por longitud
        total = len((lead or "") + (rest or ""))
        lines = max(1, int(total / (w.inches * 5.2)) + 1)
        y = y + Inches(0.32 * lines + gap)
    return y


def arrow(s, l, t, w=Inches(0.5)):
    txt(s, l, t, w, Inches(0.5), "→", 22, DIM, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def logo_mark(s, l, t, size=Inches(0.8)):
    # 2x2 de colores de severidad como identidad
    q = Emu(int(size) // 2)
    cols = [[BAJO, MEDIO], [ALTO, CRIT]]
    for r in range(2):
        for c in range(2):
            sp = rect(s, l + q*c, t + q*r, q, q, fill=cols[r][c], line=None, radius=0.0, shape=MSO_SHAPE.RECTANGLE)
    # marco redondeado encima (transparente con borde)


# ============================ SLIDES ============================

# 1. Portada
s = slide()
logo_mark(s, ML, Inches(1.5), Inches(0.9))
eyebrow(s, "Seguridad Informática · ITIZ3301 · Proyecto Integrador", t=Inches(2.7))
rich(s, ML, Inches(3.1), Inches(11.5), Inches(1.4),
     [("CyberRisk ", 58, TEXT, True, HEAD), ("360", 58, BRAND, True, HEAD)])
txt(s, ML, Inches(4.35), Inches(9.6), Inches(1.1),
    "Automatización de la gestión de riesgos cibernéticos con la metodología GRC-360, "
    "inteligencia de amenazas real y cumplimiento internacional.", 18, MUTED, spacing=1.25)
pill = rect(s, ML, Inches(5.55), Inches(6.6), Inches(0.55), fill=RGBColor(0x14,0x1f,0x38), line=RGBColor(0x2d,0x4d,0x87), radius=0.5)
shape_text(pill, "https://mrflippermen.github.io/proyecto-integrador-grc360/", 12.5,
           RGBColor(0xa9,0xc4,0xfb), bold=True, font="Consolas", align=PP_ALIGN.CENTER)
footer(s, "Integrantes del grupo · Presentación ejecutiva")

# 2. Problema
s = slide()
eyebrow(s, "El contexto")
title(s, "Gestionar el riesgo cibernético es difícil… y caro")
cw, gap = Inches(3.63), Inches(0.28)
x = ML; y = Inches(2.55); ch = Inches(2.4)
data = [("USD 4.8M", CRIT, "Costo promedio global de una brecha de datos. La mayoría de PYMES no puede absorberlo."),
        ("Hojas de cálculo", ALTO, "Sin trazabilidad, sin priorización y sin actualización: así se gestiona el riesgo aún hoy."),
        ("¿Qué primero?", MEDIO, "Sin inteligencia de amenazas no se sabe qué se está explotando ahora en el mundo real.")]
for val, col, body in data:
    c = rect(s, x, y, cw, ch)
    txt(s, x + Inches(0.28), y + Inches(0.3), cw - Inches(0.56), Inches(1.0), val, 30, col, bold=True, font=HEAD, spacing=1.0)
    txt(s, x + Inches(0.28), y + Inches(1.35), cw - Inches(0.56), Inches(0.95), body, 12.5, MUTED, spacing=1.12)
    x = x + cw + gap
txt(s, ML, Inches(5.35), Inches(11.4), Inches(1.0),
    "Impacto económico, social y global: detrás de cada activo hay datos de personas que dependen de decisiones muchas veces invisibles.",
    15, MUTED, spacing=1.2)
footer(s, "El problema")

# 3. Metodología
s = slide()
eyebrow(s, "Nuestra propuesta metodológica")
title(s, "Metodología GRC-360")
txt(s, ML, Inches(2.05), Inches(11), Inches(0.5), "Un ciclo de seis fases, automatizable y trazable, alineado con estándares internacionales.", 15, MUTED)
fases = [("1","Valoración de activos","Valoración CIA (1–5)."),
         ("2","Identificación de riesgos","Amenaza + vulnerabilidad · P × I."),
         ("3","Tratamiento","Controles ISO/IEC 27002:2022."),
         ("4","Riesgo residual","Recalculo automático."),
         ("5","Comunicación","Observaciones e informes."),
         ("6","Monitoreo","Panel, KPIs y tendencia.")]
cw = Inches(3.63); ch = Inches(1.62); gx = Inches(0.28); gy = Inches(0.2)
x0 = ML; y0 = Inches(2.6)
for idx, (n, hd, bd) in enumerate(fases):
    col = idx % 3; row = idx // 3
    x = x0 + col*(cw+gx); y = y0 + row*(ch+gy)
    card(s, x, y, cw, ch, num=n, head=hd, body=bd)
chips(s, ["ISO/IEC 27001","27002:2022","27005","ISO 31000","NIST CSF 2.0","CVSS v3.1"], ML, Inches(6.28))
footer(s, "Metodología")

# 4. Arquitectura
s = slide()
eyebrow(s, "Cómo está construido")
title(s, "Arquitectura del sistema")
layers = [("FRONTEND", "Astro (SSR) + SCSS + TypeScript · gráficos Chart.js · app cliente estática en GitHub Pages"),
          ("BACKEND / API", "Flask (API REST) · lógica GRC-360: CVSS 3.1, riesgo residual, VPR, Cyber Exposure Score, SLA"),
          ("DATOS · INTELIGENCIA", "SQLite + SQLAlchemy   ·   NVD · EPSS · CISA KEV (APIs de inteligencia en vivo)")]
y = Inches(2.55)
for lk, lt in layers:
    c = rect(s, ML, y, Inches(11.5), Inches(1.0))
    txt(s, ML + Inches(0.3), y, Inches(2.3), Inches(1.0), lk, 11.5, BRAND2, bold=True, font="Consolas", anchor=MSO_ANCHOR.MIDDLE)
    txt(s, ML + Inches(2.7), y, Inches(8.5), Inches(1.0), lt, 14, TEXT, anchor=MSO_ANCHOR.MIDDLE, spacing=1.12)
    if lk != layers[-1][0]:
        txt(s, Inches(6.4), y + Inches(1.0), Inches(0.6), Inches(0.35), "↓", 18, DIM, align=PP_ALIGN.CENTER)
    y = y + Inches(1.35)
footer(s, "Arquitectura")

# 5. Fases 1-2 con matriz
s = slide()
eyebrow(s, "Fases 1 y 2")
title(s, "De los activos a la matriz de riesgo")
bullets(s, ML, Inches(2.5), Inches(6.2), [
    ("Valoración CIA:", "cada activo se califica 1–5 en Confidencialidad, Integridad y Disponibilidad; el valor es la dimensión más crítica."),
    ("Riesgo =", "Activo + Amenaza + Vulnerabilidad."),
    ("Riesgo inherente =", "Probabilidad × Impacto (rango 1–25)."),
], size=15, gap=0.22)
# panel matriz
panel = rect(s, Inches(7.55), Inches(2.4), Inches(4.9), Inches(3.85), fill=SURF)
txt(s, Inches(7.85), Inches(2.62), Inches(4.4), Inches(0.3), "MATRIZ DE RIESGO 5×5", 10.5, MUTED, bold=True, font=HEAD)
cell = Inches(0.5); g = Inches(0.075); mx = Inches(7.9); my = Inches(3.02)
for f in range(5):
    for c in range(5):
        imp = 5 - f; prob = c + 1; sc = imp * prob
        col = BAJO if sc <= 4 else MEDIO if sc <= 9 else ALTO if sc <= 14 else CRIT
        cc = rect(s, mx + c*(cell+g), my + f*(cell+g), cell, cell, fill=col, line=None, radius=0.14)
        shape_text(cc, str(sc), 11, WHITE, bold=True, font=HEAD)
leg = ["Bajo", "Medio", "Alto", "Crítico"]; lc = [BAJO, MEDIO, ALTO, CRIT]; lx = Inches(7.9)
for i, lab in enumerate(leg):
    rect(s, lx, Inches(5.94), Inches(0.16), Inches(0.16), fill=lc[i], line=None, radius=0.1, shape=MSO_SHAPE.RECTANGLE)
    txt(s, lx + Inches(0.22), Inches(5.9), Inches(0.9), Inches(0.3), lab, 10.5, MUTED, anchor=MSO_ANCHOR.MIDDLE)
    lx = lx + Inches(1.05)
footer(s, "Valoración e identificación")

# 6. Fases 3-4
s = slide()
eyebrow(s, "Fases 3 y 4")
title(s, "Tratamiento y riesgo residual automático")
chips(s, ["Mitigar", "Transferir", "Aceptar", "Evitar"], ML, Inches(2.15))
bullets(s, ML, Inches(2.95), Inches(11.4), [
    ("", "Cada tratamiento se asocia a un control ISO/IEC 27002:2022, un responsable y una eficacia estimada."),
    ("Riesgo residual = Inherente × (1 − eficacia)", "de los controles, combinada de forma multiplicativa (rendimientos decrecientes)."),
    ("", "Solo los controles implementados reducen el riesgo, y nunca llega a cero (máximo 90 % de reducción)."),
    ("", "El residual, el VPR y el Cyber Exposure Score se recalculan solos con cada cambio."),
], size=15, gap=0.26)
footer(s, "Tratamiento y residual")

# 7. CTI
s = slide()
eyebrow(s, "Inteligencia de amenazas")
title(s, "Priorizar por la amenaza real, no solo la teórica")
tri = [("CVSS 3.1", BRAND2, "Severidad técnica oficial (NVD / FIRST.org)."),
       ("EPSS", MEDIO, "Probabilidad de explotación en los próximos 30 días."),
       ("CISA KEV", CRIT, "¿Se está explotando activamente en el mundo real?")]
cw = Inches(3.63); x = ML; y = Inches(2.5)
for hd, col, bd in tri:
    card(s, x, y, cw, Inches(1.9), head=hd, body=bd, head_color=col)
    x = x + cw + Inches(0.28)
big = rect(s, ML, Inches(4.65), Inches(11.5), Inches(1.4), fill=SURF)
txt(s, ML + Inches(0.35), Inches(4.85), Inches(7.5), Inches(1.0),
    "VPR — Vulnerability Priority Rating", 17, TEXT, bold=True, font=HEAD)
txt(s, ML + Inches(0.35), Inches(5.35), Inches(7.5), Inches(0.6),
    "Combina CVSS + EPSS + KEV en una sola prioridad 0–10, al estilo Tenable.", 13, MUTED)
txt(s, Inches(9.0), Inches(4.9), Inches(3.1), Inches(0.4), "Ej. Log4Shell (CVE-2021-44228)", 11, DIM, align=PP_ALIGN.RIGHT)
txt(s, Inches(9.0), Inches(5.25), Inches(3.1), Inches(0.6), "VPR 10 · Crítica", 22, BRAND2, bold=True, font=HEAD, align=PP_ALIGN.RIGHT)
footer(s, "Threat Intelligence")

# 8. Automatización
s = slide()
eyebrow(s, "Diferenciador clave")
rich(s, ML, Inches(1.25), Inches(11.5), Inches(1.0),
     [("Valoración y mitigación ", 34, TEXT, True, HEAD), ("automáticas", 34, BRAND, True, HEAD)])
card(s, ML, Inches(2.6), Inches(5.6), Inches(2.2), head="Valoración automática",
     body="Al elegir activo + amenaza + vulnerabilidad, el sistema calcula solo la probabilidad "
          "(desde EPSS/KEV) y el impacto (criticidad del activo + CVSS).")
card(s, Inches(6.78), Inches(2.6), Inches(5.62), Inches(2.2), head="Mitigación automática",
     body="Un motor de reglas recomienda los controles ISO 27002 adecuados, proyecta el riesgo "
          "residual y los aplica con un clic — individual o masivo.")
txt(s, ML, Inches(5.15), Inches(11.4), Inches(1.0),
    "El analista decide; el sistema hace los cálculos y propone. La responsabilidad profesional no se "
    "automatiza: todo queda en la bitácora de auditoría.", 15, MUTED, spacing=1.2)
footer(s, "Automatización")

# 9. Cumplimiento
s = slide()
eyebrow(s, "Cumplimiento")
title(s, "Alineado a marcos internacionales")
card(s, ML, Inches(2.6), Inches(5.6), Inches(2.2), head="ISO/IEC 27001 · 27002:2022",
     body="Cobertura del programa frente a los 93 controles del Anexo A, agrupados por tema: "
          "Organizacional, Personas, Físico y Tecnológico.")
card(s, Inches(6.78), Inches(2.6), Inches(5.62), Inches(2.2), head="NIST CSF 2.0",
     body="Distribución de controles en las 6 funciones del marco, revelando dónde reforzar la defensa.")
chips(s, ["Govern", "Identify", "Protect", "Detect", "Respond", "Recover"], ML, Inches(5.2), on=False)
footer(s, "Cumplimiento")

# 10. Monitoreo
s = slide()
eyebrow(s, "Fase 6")
title(s, "Panel de monitoreo")
stats = [("Cyber Exposure Score", "423 / 1000", MEDIO, "Exposición agregada de la organización."),
         ("Apetito de riesgo", "6", CRIT, "Riesgos que exceden el umbral tolerado."),
         ("SLA de remediación", "6", ALTO, "Riesgos fuera de su ventana de tratamiento.")]
cw = Inches(3.63); x = ML; y = Inches(2.55)
for lab, val, col, bd in stats:
    rect(s, x, y, cw, Inches(2.3))
    txt(s, x + Inches(0.28), y + Inches(0.26), cw - Inches(0.56), Inches(0.4), lab.upper(), 10.5, MUTED, bold=True, font=HEAD)
    txt(s, x + Inches(0.28), y + Inches(0.72), cw - Inches(0.56), Inches(0.8), val, 30, col, bold=True, font=HEAD)
    txt(s, x + Inches(0.28), y + Inches(1.55), cw - Inches(0.56), Inches(0.7), bd, 12.5, MUTED, spacing=1.1)
    x = x + cw + Inches(0.28)
txt(s, ML, Inches(5.2), Inches(11.4), Inches(0.8),
    "Más matriz de calor 5×5, distribución inherente vs. residual y tendencia de exposición histórica.", 15, MUTED)
footer(s, "Monitoreo")

# 11. Flujo
s = slide()
eyebrow(s, "Flujo de trabajo de la herramienta")
title(s, "De un archivo Excel a decisiones informadas")
steps = [("1 · CARGAR", "Importar", "Activos y riesgos desde Excel/CSV."),
         ("2 · VALORAR", "Auto P × I", "El sistema sugiere probabilidad e impacto."),
         ("3 · TRATAR", "Mitigación", "Controles ISO recomendados y aplicados."),
         ("4 · MONITOREAR", "Panel", "Residual, exposición y cumplimiento en vivo.")]
bw = Inches(2.62); gap = Inches(0.35); x = ML; y = Inches(2.9); bh = Inches(1.9)
for i, (k, t2, d) in enumerate(steps):
    rect(s, x, y, bw, bh)
    txt(s, x + Inches(0.22), y + Inches(0.24), bw - Inches(0.44), Inches(0.35), k, 10.5, BRAND2, bold=True, font=HEAD)
    txt(s, x + Inches(0.22), y + Inches(0.66), bw - Inches(0.44), Inches(0.45), t2, 16, TEXT, bold=True, font=HEAD)
    txt(s, x + Inches(0.22), y + Inches(1.12), bw - Inches(0.44), Inches(0.7), d, 12, MUTED, spacing=1.1)
    if i < 3:
        txt(s, x + bw, y, gap, bh, "→", 20, DIM, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x = x + bw + gap
txt(s, ML, Inches(5.35), Inches(11.4), Inches(0.7),
    "Todo el ciclo GRC-360 en un solo flujo, intuitivo y sin instalar nada.", 15, MUTED)
footer(s, "Flujo de trabajo")

# 12. Ética RC4
s = slide()
eyebrow(s, "Responsabilidad ética y profesional · RC4")
rich(s, ML, Inches(1.25), Inches(11.5), Inches(1.0),
     [("La decisión de diseño más importante es ", 32, TEXT, True, HEAD), ("ética", 32, BRAND, True, HEAD)])
bullets(s, ML, Inches(2.7), Inches(11.4), [
    ("Honestidad sobre los límites:", "el riesgo residual nunca es cero y un tablero en verde no significa “estar seguro”."),
    ("Rendición de cuentas:", "la automatización no exime al profesional; todo queda trazado en la bitácora de auditoría."),
    ("Impacto en las personas:", "proteger datos es un acto de respeto, no un mero requisito de cumplimiento."),
    ("Conocimiento de doble uso:", "la inteligencia de amenazas se usa con una finalidad exclusivamente defensiva."),
], size=15, gap=0.3)
footer(s, "Ética profesional")

# 13. Demo
s = slide()
eyebrow(s, "Demostración en vivo", t=Inches(2.3))
rich(s, ML, Inches(2.75), Inches(11.5), Inches(1.6),
     [("Pruébelo ahora — ", 46, TEXT, True, HEAD), ("sin instalar nada", 46, BRAND, True, HEAD)])
pill = rect(s, ML, Inches(4.5), Inches(7.4), Inches(0.62), fill=RGBColor(0x14,0x1f,0x38), line=RGBColor(0x2d,0x4d,0x87), radius=0.5)
shape_text(pill, "https://mrflippermen.github.io/proyecto-integrador-grc360/", 14,
           RGBColor(0xa9,0xc4,0xfb), bold=True, font="Consolas", align=PP_ALIGN.CENTER)
txt(s, ML, Inches(5.4), Inches(10.5), Inches(0.9),
    "Abra el enlace en cualquier navegador. Caso de ejemplo ya cargado · carga de Excel · inteligencia de amenazas · mitigación automática.",
    15, MUTED, spacing=1.2)
footer(s, "Demo")

# 14. Conclusiones
s = slide()
eyebrow(s, "Conclusiones")
title(s, "Tecnología de seguridad potente… y responsable")
bullets(s, ML, Inches(2.55), Inches(11.4), [
    ("", "Automatizamos íntegramente la metodología GRC-360, de la valoración al monitoreo continuo."),
    ("", "Alcanzamos estándares internacionales (ISO, NIST, CVSS, MITRE) con inteligencia de amenazas real."),
    ("", "La valoración y la mitigación se calculan automáticamente, ayudando a decidir mejor y más rápido."),
], size=15, gap=0.28)
txt(s, ML, Inches(4.9), Inches(11), Inches(0.8), "Gracias. ¿Preguntas?", 24, WHITE, bold=True, font=HEAD)
footer(s, "ITIZ3301 · Proyecto Integrador")

OUT = "/home/flippermen/Downloads/proyecto-integrador-grc360/presentacion/CyberRisk360-Presentacion.pptx"
prs.save(OUT)
print("Guardado:", OUT)
print("Slides:", len(prs.slides._sldIdLst))
